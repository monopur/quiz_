import os
import sys
import shutil
import re
import sqlite3
from pptx import Presentation

DB_PATH = os.path.join("..", "db", "quiz.db")  # Veritabanı yolu
MEDIA_DIR = os.path.join("..", "static", "uploads", "pptx_media")  # Medya dosyalarının kaydedileceği klasör

def extract_notes(slide):
    notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
    # Notlar kısmından doğru cevap, kategori ve zorluk
    correct, cat, diff = "", "", ""
    for line in notes.split("\n"):
        l = line.strip()
        if re.match(r"(?i)doğru.*:", l):
            correct = l.split(":", 1)[1].strip()
        elif re.match(r"(?i)kat.*:", l):
            cat = l.split(":", 1)[1].strip()
        elif re.match(r"(?i)zorluk.*:", l):
            diff = l.split(":", 1)[1].strip()
    return correct, cat, diff

def get_or_create_category(db, cat_name):
    cat_row = db.execute("SELECT id FROM categories WHERE name=?", (cat_name,)).fetchone()
    if cat_row:
        return cat_row[0]
    db.execute("INSERT INTO categories (name) VALUES (?)", (cat_name,))
    db.commit()
    return db.execute("SELECT id FROM categories WHERE name=?", (cat_name,)).fetchone()[0]

def insert_question(db, data):
    db.execute(
        """INSERT INTO questions (category_id, question_text, image_path, audio_path, video_path, option_a, option_b, option_c, option_d, correct_option, difficulty, points, duration)
           VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)""",
        data
    )
    db.commit()

def main(pptx_path):
    prs = Presentation(pptx_path)
    if os.path.exists(MEDIA_DIR):
        shutil.rmtree(MEDIA_DIR)
    os.makedirs(MEDIA_DIR, exist_ok=True)
    db = sqlite3.connect(DB_PATH)
    for idx, slide in enumerate(prs.slides, 1):
        q_text = ""
        img_path, audio_path, video_path = None, None, None
        # Metin kutularını bul
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip() and not q_text:
                q_text = shape.text.strip()
            if shape.shape_type == 13:  # Picture
                img_file = f"slide{idx}_img.png"
                with open(os.path.join(MEDIA_DIR, img_file), "wb") as img_f:
                    img_f.write(shape.image.blob)
                img_path = f"pptx_media/{img_file}"
        # Ses/video dosyaları (media_part)
        for rel in slide.part.rels.values():
            if "audio" in rel.target_ref:
                audio_file = f"slide{idx}_audio.mp3"
                target = rel.target_part.blob
                with open(os.path.join(MEDIA_DIR, audio_file), "wb") as aud_f:
                    aud_f.write(target)
                audio_path = f"pptx_media/{audio_file}"
            if "video" in rel.target_ref:
                video_file = f"slide{idx}_video.mp4"
                target = rel.target_part.blob
                with open(os.path.join(MEDIA_DIR, video_file), "wb") as vid_f:
                    vid_f.write(target)
                video_path = f"pptx_media/{video_file}"
        correct, cat, diff = extract_notes(slide)
        if not cat:
            cat = "Genel"
        cat_id = get_or_create_category(db, cat)
        # Alternatifler için Notlarda ya da scriptte placeholder (geliştirilebilir)
        option_a = "A seçeneği"
        option_b = "B seçeneği"
        option_c = "C seçeneği"
        option_d = "D seçeneği"
        points = 10
        duration = 20
        data = (
            cat_id, q_text, img_path, audio_path, video_path,
            option_a, option_b, option_c, option_d,
            correct, diff or "orta", points, duration
        )
        insert_question(db, data)
        print(f"Soru eklendi: {q_text[:30]}... [Kategori: {cat}, Zorluk: {diff}]")
    db.close()
    print(f"Tüm sorular ve medya '{DB_PATH}' veritabanına ve '{MEDIA_DIR}/' klasörüne eklendi.")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Kullanım: python pptx2quiz_db.py sunu.pptx")
        sys.exit(1)
    main(sys.argv[1])
