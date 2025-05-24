import os
import sys
import shutil
import re
from pptx import Presentation

OUTPUT_CSV = "quiz_import.csv"
MEDIA_DIR = "quiz_media"

def extract_notes(slide):
    notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
    # Doğru cevap, kategori ve zorluk bilgisini satır satır ayır
    correct, cat, diff = "", "", ""
    for line in notes.split("\n"):
        if re.match(r"(?i)doğru.*:", line):
            correct = line.split(":", 1)[1].strip()
        elif re.match(r"(?i)kat.*:", line):
            cat = line.split(":", 1)[1].strip()
        elif re.match(r"(?i)zorluk.*:", line):
            diff = line.split(":", 1)[1].strip()
    return correct, cat, diff

def main(pptx_path):
    prs = Presentation(pptx_path)
    if os.path.exists(MEDIA_DIR):
        shutil.rmtree(MEDIA_DIR)
    os.makedirs(MEDIA_DIR, exist_ok=True)
    with open(OUTPUT_CSV, "w", encoding="utf-8") as f:
        f.write("category_id,question_text,image_path,audio_path,video_path,option_a,option_b,option_c,option_d,correct_option,difficulty,points,duration\n")
        for idx, slide in enumerate(prs.slides, 1):
            q_text = ""
            img_path, audio_path, video_path = "", "", ""
            # Metin kutularını bul
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    q_text += shape.text.strip() + " "
                # Resim dosyalarını kaydet
                if shape.shape_type == 13:  # Picture
                    img_file = f"slide{idx}_img.png"
                    with open(os.path.join(MEDIA_DIR, img_file), "wb") as img_f:
                        img_f.write(shape.image.blob)
                    img_path = img_file
            # Ses/video dosyaları (media_part)
            for rel in slide.part.rels.values():
                if "audio" in rel.target_ref:
                    audio_file = f"slide{idx}_audio.mp3"
                    target = rel.target_part.blob
                    with open(os.path.join(MEDIA_DIR, audio_file), "wb") as aud_f:
                        aud_f.write(target)
                    audio_path = audio_file
                if "video" in rel.target_ref:
                    video_file = f"slide{idx}_video.mp4"
                    target = rel.target_part.blob
                    with open(os.path.join(MEDIA_DIR, video_file), "wb") as vid_f:
                        vid_f.write(target)
                    video_path = video_file
            correct, cat, diff = extract_notes(slide)
            # Alternatifler için placeholder (A/B/C/D)
            option_a = "A seçeneği"
            option_b = "B seçeneği"
            option_c = "C seçeneği"
            option_d = "D seçeneği"
            points = 10
            duration = 20
            f.write(f"{cat},{q_text.strip()},{img_path},{audio_path},{video_path},{option_a},{option_b},{option_c},{option_d},{correct},{diff},{points},{duration}\n")
    print(f"CSV ve medya dosyaları '{OUTPUT_CSV}' ve '{MEDIA_DIR}/' klasörüne yazıldı.")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Kullanım: python pptx2quiz.py dosya.pptx")
        sys.exit(1)
    main(sys.argv[1])
