import os
import shutil
import re
import sqlite3
from pptx import Presentation

DB_PATH = os.path.join(os.path.dirname(__file__), "..", "db", "quiz.db")
MEDIA_DIR = os.path.join(os.path.dirname(__file__), "..", "static", "uploads", "pptx_media")

def extract_notes(slide):
    notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
    correct, cat, diff = "", "", ""
    for line in notes.split("\n"):
        l = line.strip()
        if re.match(r"(?i)doğru.*:", l):
            correct = l.split(":", 1)[1].strip()
        elif re.match(r"(?i)kat.*:", l):
            cat = l.split(":", 1)[1].strip()
        elif re.match(r"(?i)zorluk.*:", l):
            diff = l.split(":", 1)[1].strip()
    return correct, cat, diff

def get_or_create_category(db, cat_name):
    if not cat_name:
        cat_name = "Genel"
    cat_row = db.execute("SELECT id FROM categories WHERE name=?", (cat_name,)).fetchone()
    if cat_row:
        return cat_row[0]
    db.execute("INSERT INTO categories (name) VALUES (?)", (cat_name,))
    db.commit()
    return db.execute("SELECT id FROM categories WHERE name=?", (cat_name,)).fetchone()[0]

def insert_question(db, data):
    db.execute(
        """INSERT INTO questions (category_id, question_text, image_path, audio_path, video_path, option_a, option_b, option_c, option_d, correct_option, difficulty, points, duration)
           VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)""",
        data
    )
    db.commit()

def pptx_to_db_with_result(pptx_path):
    errors = []
    added = 0
    try:
        prs = Presentation(pptx_path)
        if os.path.exists(MEDIA_DIR):
            shutil.rmtree(MEDIA_DIR)
        os.makedirs(MEDIA_DIR, exist_ok=True)
        db = sqlite3.connect(DB_PATH)
        for idx, slide in enumerate(prs.slides, 1):
            try:
                q_text = ""
                img_path, audio_path, video_path = None, None, None
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip() and not q_text:
                        q_text = shape.text.strip()
                    if shape.shape_type == 13:  # Picture
                        img_file = f"slide{idx}_img.png"
                        with open(os.path.join(MEDIA_DIR, img_file), "wb") as img_f:
                            img_f.write(shape.image.blob)
                        img_path = f"pptx_media/{img_file}"
                for rel in slide.part.rels.values():
                    if "audio" in rel.target_ref:
                        audio_file = f"slide{idx}_audio.mp3"
                        target = rel.target_part.blob
                        with open(os.path.join(MEDIA_DIR, audio_file), "wb") as aud_f:
                            aud_f.write(target)
                        audio_path = f"pptx_media/{audio_file}"
                    if "video" in rel.target_ref:
                        video_file = f"slide{idx}_video.mp4"
                        target = rel.target_part.blob
                        with open(os.path.join(MEDIA_DIR, video_file), "wb") as vid_f:
                            vid_f.write(target)
                        video_path = f"pptx_media/{video_file}"
                correct, cat, diff = extract_notes(slide)
                if not q_text:
                    errors.append(f"{idx}. slaytta soru metni yok.")
                    continue
                if not correct:
                    errors.append(f"{idx}. slaytta 'Doğru' notu eksik.")
                    continue
                cat_id = get_or_create_category(db, cat)
                # Alternatifler için şimdilik placeholder
                option_a = "A seçeneği"
                option_b = "B seçeneği"
                option_c = "C seçeneği"
                option_d = "D seçeneği"
                points = 10
                duration = 20
                data = (
                    cat_id, q_text, img_path, audio_path, video_path,
                    option_a, option_b, option_c, option_d,
                    correct, diff or "orta", points, duration
                )
                insert_question(db, data)
                added += 1
            except Exception as slideerr:
                errors.append(f"{idx}. slaytta hata: {slideerr}")
        db.close()
        if added > 0 and not errors:
            return {"success": True, "message": f"Tüm {added} soru başarıyla aktarıldı.", "errors": []}
        elif added > 0:
            return {"success": True, "message": f"{added} soru aktarıldı, bazı slaytlar eksik/hatalı.", "errors": errors}
        else:
            return {"success": False, "message": "Hiç soru eklenemedi!", "errors": errors}
    except Exception as e:
        return {"success": False, "message": f"İşlem yapılamadı: {e}", "errors": [str(e)]}
